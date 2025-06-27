import os
import base64
from io import BytesIO
import filetype
from pdf2image import convert_from_path
from datetime import datetime
import hashlib

# New imports for PPT and DOC support
from pptx import Presentation
from docx import Document
import docx2txt
from langchain_community.document_loaders import UnstructuredPowerPointLoader, Docx2txtLoader
from langchain_community.document_loaders import UnstructuredWordDocumentLoader

from langchain_core.messages import HumanMessage
from langchain_community.document_loaders import PyPDFLoader
from langchain_chroma import Chroma
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from .prompts import image_prompt  # Make sure this exists

# --- Utility Functions ---

def get_filename(file_path):
    return os.path.splitext(os.path.basename(file_path))[0]

def get_file_hash(file_path):
    """Generate SHA-256 hash of file for duplicate detection"""
    try:
        with open(file_path, 'rb') as f:
            return hashlib.sha256(f.read()).hexdigest()
    except:
        return None

def get_file_size(file_path):
    """Get file size in bytes"""
    try:
        return os.path.getsize(file_path)
    except:
        return None

def create_base_metadata(file_path, company_name, file_type):
    """Create base metadata common to all file types"""
    timestamp = datetime.now().isoformat()
    
    metadata = {
        'company_name': company_name,
        'file_type': file_type,
        'filename': os.path.basename(file_path) if isinstance(file_path, str) else 'uploaded_image',
        'file_path': file_path if isinstance(file_path, str) else None,
        'processed_timestamp': timestamp,
        'processing_date': datetime.now().strftime('%Y-%m-%d'),
        'processing_time': datetime.now().strftime('%H:%M:%S'),
        'chunk_strategy': 'recursive_character_text_splitter',
        'embedding_model': 'huggingface_default'
    }
    
    # Add file-specific metadata if it's a file path
    if isinstance(file_path, str) and os.path.exists(file_path):
        metadata.update({
            'file_hash': get_file_hash(file_path),
            'file_size_bytes': get_file_size(file_path),
            'file_extension': os.path.splitext(file_path)[1].lower()
        })
    
    return metadata

def file_router(file):
    """Enhanced file router with PPT and DOC support"""
    try:
        # Get file extension first
        file_extension = os.path.splitext(file)[1].lower()
        
        # Handle specific extensions
        if file_extension in ['.ppt', '.pptx']:
            return 'powerpoint'
        elif file_extension in ['.doc', '.docx']:
            return 'word_document'
        
        # Use filetype for other formats
        kind = filetype.guess(file)
        if kind is None:
            return "Unknown"

        file_type = kind.mime

        if file_type.startswith("image/"):
            return 'imagesingle'

        # Check if it's a PDF with text or images
        if file_type == 'application/pdf' or file_extension == '.pdf':
            loader = PyPDFLoader(file)
            docs = loader.load()

            if not len(docs[0].page_content.strip()):
                return 'imagepdf'
            else:
                return 'pdf'
                
        return 'pdf'  # Default fallback
        
    except Exception as e:
        print(f"Error in file_router: {e}")
        return 'pdf'  # Default fallback

def encode_image(image) -> str:
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return base64.b64encode(buffer.getvalue()).decode("utf-8")

# --- LLM Setup ---

model = ChatGoogleGenerativeAI(model='gemini-2.0-flash')

def image_summarize(model, base64_image: str, prompt: str) -> str:
    msg = model.invoke([
        HumanMessage(
            content=[
                {"type": "text", "text": prompt},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{base64_image}"},
                },
            ]
        )
    ])
    return msg.content

# --- Image Handlers ---

def image_handler(image):
    base64_img = encode_image(image)
    summary = image_summarize(model, base64_img, prompt=image_prompt)
    with open('example.txt', 'w') as f:
        f.write(summary)
    return summary

def image_handler_append(image):
    base64_img = encode_image(image)
    summary = image_summarize(model, base64_img, prompt=image_prompt)
    with open('example.txt', 'a') as f:
        f.write(summary + '\n')
    return summary

# --- PowerPoint Handler ---

def extract_ppt_content(filepath: str):
    """Extract text content from PowerPoint files"""
    try:
        prs = Presentation(filepath)
        full_text = []
        slide_count = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"=== Slide {slide_num} ===\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
                    
                # Handle tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text += " | ".join(row_text) + "\n"
            
            if slide_text.strip() != f"=== Slide {slide_num} ===":
                full_text.append(slide_text)
                slide_count += 1
        
        return "\n\n".join(full_text), slide_count
        
    except Exception as e:
        print(f"Error extracting PowerPoint content with python-pptx: {e}")
        # Fallback to langchain loader
        try:
            loader = UnstructuredPowerPointLoader(filepath)
            docs = loader.load()
            content = "\n\n".join([doc.page_content for doc in docs])
            return content, len(docs)
        except Exception as fallback_error:
            print(f"Fallback PowerPoint loader failed: {fallback_error}")
            return f"Error processing PowerPoint file: {str(e)}", 0

# --- Word Document Handler ---

def extract_word_content(filepath: str):
    """Extract text content from Word documents"""
    try:
        file_extension = os.path.splitext(filepath)[1].lower()
        
        if file_extension == '.docx':
            # Use python-docx for .docx files
            doc = Document(filepath)
            full_text = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        full_text.append(" | ".join(row_text))
            
            content = "\n\n".join(full_text)
            return content, len(doc.paragraphs)
            
        elif file_extension == '.doc':
            # Use python-docx2txt for .doc files
            content = docx2txt.process(filepath)
            return content, len(content.split('\n'))
            
    except Exception as e:
        print(f"Error extracting Word content: {e}")
        # Fallback to langchain loaders
        try:
            if filepath.endswith('.docx'):
                loader = Docx2txtLoader(filepath)
            else:
                loader = UnstructuredWordDocumentLoader(filepath)
            
            docs = loader.load()
            content = "\n\n".join([doc.page_content for doc in docs])
            return content, len(docs)
            
        except Exception as fallback_error:
            print(f"Fallback Word loader failed: {fallback_error}")
            return f"Error processing Word document: {str(e)}", 0

# --- Enhanced Vectorization Functions ---

def vectorize_text(text: str, company_name: str, filename: str = "text_input", base_metadata: dict = None):
    """Vectorize text content with metadata"""
    try:
        splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=50)
        docs = splitter.split_text(text)
        
        # Create persist directory
        persist_directory = os.path.join("chroma_store", company_name, filename)
        os.makedirs(persist_directory, exist_ok=True)
        
        # Create collection name (sanitize company name)
        collection_name = f"{company_name}_{filename}".replace(" ", "_").replace("-", "_").lower()
        
        # Create metadata for each chunk
        metadatas = []
        for i, chunk in enumerate(docs):
            chunk_metadata = base_metadata.copy() if base_metadata else {}
            chunk_metadata.update({
                'chunk_index': i,
                'chunk_size': len(chunk),
                'total_chunks': len(docs),
                'content_type': 'text',
                'source_document': filename
            })
            metadatas.append(chunk_metadata)
        
        vectorstore = Chroma.from_texts(
            texts=docs,
            embedding=HuggingFaceEmbeddings(),
            metadatas=metadatas,
            persist_directory=persist_directory,
            collection_name=collection_name
        )
        return vectorstore
        
    except Exception as e:
        print(f"Error in vectorize_text: {e}")
        # Fallback to in-memory store
        metadatas = [{'error': str(e), 'fallback': True} for _ in docs]
        vectorstore = Chroma.from_texts(
            texts=docs,
            embedding=HuggingFaceEmbeddings(),
            metadatas=metadatas,
            collection_name=f"fallback_{collection_name}"
        )
        return vectorstore

def vectorize_powerpoint(filepath: str, company_name: str):
    """Vectorize PowerPoint presentations"""
    try:
        content, slide_count = extract_ppt_content(filepath)
        filename = get_filename(filepath)
        
        # Create base metadata
        base_metadata = create_base_metadata(filepath, company_name, 'powerpoint')
        base_metadata.update({
            'total_slides': slide_count,
            'content_source': 'powerpoint_extraction',
            'extraction_method': 'python_pptx_with_langchain_fallback',
            'supports_tables': True,
            'supports_shapes': True
        })
        
        return vectorize_text(content, company_name, filename, base_metadata)
        
    except Exception as e:
        print(f"Error in vectorize_powerpoint: {e}")
        error_metadata = {
            'error': str(e), 
            'file_type': 'powerpoint',
            'extraction_failed': True
        }
        return vectorize_text("Error processing PowerPoint file", company_name, "error_ppt", error_metadata)

def vectorize_word_document(filepath: str, company_name: str):
    """Vectorize Word documents"""
    try:
        content, paragraph_count = extract_word_content(filepath)
        filename = get_filename(filepath)
        
        # Create base metadata
        base_metadata = create_base_metadata(filepath, company_name, 'word_document')
        base_metadata.update({
            'paragraph_count': paragraph_count,
            'content_source': 'word_extraction',
            'extraction_method': 'python_docx_with_langchain_fallback',
            'supports_tables': True,
            'supports_formatting': True
        })
        
        return vectorize_text(content, company_name, filename, base_metadata)
        
    except Exception as e:
        print(f"Error in vectorize_word_document: {e}")
        error_metadata = {
            'error': str(e), 
            'file_type': 'word_document',
            'extraction_failed': True
        }
        return vectorize_text("Error processing Word document", company_name, "error_doc", error_metadata)

def vectorize_single_image(image, company_name: str):
    """Vectorize single images"""
    try:
        # Create base metadata for image
        base_metadata = create_base_metadata(image, company_name, 'single_image')
        base_metadata.update({
            'content_source': 'ai_image_summary',
            'ai_model_used': 'gemini-2.0-flash',
            'processing_method': 'image_to_text_summary'
        })
        
        summary = image_handler(image)
        filename = "image_single"
        return vectorize_text(summary, company_name, filename, base_metadata)
    except Exception as e:
        print(f"Error in vectorize_single_image: {e}")
        error_metadata = {'error': str(e), 'file_type': 'single_image'}
        return vectorize_text("Error processing image", company_name, "error_image", error_metadata)

def vectorize_multiple_images(image_path: str, company_name: str):
    """Vectorize PDF with images"""
    try:
        images = convert_from_path(image_path)
        filename = get_filename(image_path)
        
        # Create base metadata
        base_metadata = create_base_metadata(image_path, company_name, 'pdf_images')
        base_metadata.update({
            'total_pages': len(images),
            'content_source': 'ai_image_summary',
            'ai_model_used': 'gemini-2.0-flash',
            'processing_method': 'pdf_to_images_to_text',
            'conversion_tool': 'pdf2image'
        })
        
        summary = ''
        for i, image in enumerate(images):
            if i == 0:
                summary = image_handler(image)
            else:
                summary += image_handler_append(image)

        return vectorize_text(summary, company_name, filename, base_metadata)
    except Exception as e:
        print(f"Error in vectorize_multiple_images: {e}")
        error_metadata = {'error': str(e), 'file_type': 'pdf_images'}
        return vectorize_text("Error processing PDF images", company_name, "error_pdf_images", error_metadata)

def vectorize_docs(filepath: str, company_name: str):
    """Vectorize PDF documents"""
    try:
        loader = PyPDFLoader(filepath)
        docs = loader.load()
        splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=80)
        chunks = splitter.split_documents(docs)
        filename = get_filename(filepath)
        
        # Create base metadata
        base_metadata = create_base_metadata(filepath, company_name, 'pdf_document')
        base_metadata.update({
            'total_pages': len(docs),
            'total_chunks_created': len(chunks),
            'chunk_size': 600,
            'chunk_overlap': 80,
            'content_source': 'direct_pdf_text',
            'loader_used': 'PyPDFLoader'
        })
        
        # Create persist directory
        persist_directory = os.path.join("chroma_store", company_name, filename)
        os.makedirs(persist_directory, exist_ok=True)
        
        # Create collection name (sanitize)
        collection_name = f"{company_name}_{filename}".replace(" ", "_").replace("-", "_").lower()
        
        # Add metadata to each chunk
        for i, chunk in enumerate(chunks):
            chunk.metadata.update(base_metadata)
            chunk.metadata.update({
                'chunk_index': i,
                'page_number': chunk.metadata.get('page', 'unknown'),
                'chunk_char_count': len(chunk.page_content)
            })
        
        vectorstore = Chroma.from_documents(
            documents=chunks,
            embedding=HuggingFaceEmbeddings(),
            persist_directory=persist_directory,
            collection_name=collection_name
        )
        return vectorstore
        
    except Exception as e:
        print(f"Error in vectorize_docs: {e}")
        # Fallback to in-memory store
        try:
            loader = PyPDFLoader(filepath)
            docs = loader.load()
            splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=80)
            chunks = splitter.split_documents(docs)
            
            # Add error metadata to fallback
            error_metadata = {'error': str(e), 'fallback': True, 'file_type': 'pdf_document'}
            for chunk in chunks:
                chunk.metadata.update(error_metadata)
            
            vectorstore = Chroma.from_documents(
                documents=chunks,
                embedding=HuggingFaceEmbeddings(),
                collection_name=f"fallback_{company_name}_{filename}".replace(" ", "_").lower()
            )
            return vectorstore
        except Exception as fallback_error:
            print(f"Fallback also failed: {fallback_error}")
            # Return minimal vectorstore
            return Chroma.from_texts(
                texts=["Error loading document"],
                embedding=HuggingFaceEmbeddings(),
                metadatas=[{'error': str(fallback_error), 'critical_failure': True}],
                collection_name="error_fallback"
            )

# --- Entry Point for Routing ---

def vectorize(filepath: str, company_name: str):
    """Main vectorization function with enhanced file type support"""
    try:
        file_type = file_router(filepath)
        print(f"Detected file type: {file_type}")

        if file_type == 'imagesingle':
            return vectorize_single_image(filepath, company_name)
        elif file_type == 'imagepdf':
            return vectorize_multiple_images(filepath, company_name)
        elif file_type == 'powerpoint':
            return vectorize_powerpoint(filepath, company_name)
        elif file_type == 'word_document':
            return vectorize_word_document(filepath, company_name)
        else:
            return vectorize_docs(filepath, company_name)
            
    except Exception as e:
        print(f"Error in vectorize main function: {e}")
        # Ultimate fallback with comprehensive error metadata
        error_metadata = {
            'error': str(e),
            'critical_failure': True,
            'processed_timestamp': datetime.now().isoformat(),
            'company_name': company_name,
            'attempted_file': filepath
        }
        return Chroma.from_texts(
            texts=[f"Error processing file: {str(e)}"],
            embedding=HuggingFaceEmbeddings(),
            metadatas=[error_metadata],
            collection_name="ultimate_fallback"
        )

# --- Utility Functions for Metadata Queries ---

def search_by_metadata(vectorstore, metadata_filter: dict, query: str = None, k: int = 5):
    """Search documents using metadata filters"""
    try:
        if query:
            # Similarity search with metadata filter
            results = vectorstore.similarity_search(
                query=query, 
                k=k, 
                filter=metadata_filter
            )
        else:
            # Get all documents matching metadata filter
            results = vectorstore.get(where=metadata_filter, limit=k)
        return results
    except Exception as e:
        print(f"Error in metadata search: {e}")
        return []

def get_document_metadata_summary(vectorstore):
    """Get a summary of all metadata in the vectorstore"""
    try:
        # This would need to be implemented based on your specific Chroma setup
        # You might need to query the underlying collection directly
        collection = vectorstore._collection
        all_data = collection.get()
        
        if all_data and 'metadatas' in all_data:
            return {
                'total_documents': len(all_data['metadatas']),
                'unique_companies': set(meta.get('company_name') for meta in all_data['metadatas'] if meta.get('company_name')),
                'file_types': set(meta.get('file_type') for meta in all_data['metadatas'] if meta.get('file_type')),
                'processing_dates': set(meta.get('processing_date') for meta in all_data['metadatas'] if meta.get('processing_date'))
            }
    except Exception as e:
        print(f"Error getting metadata summary: {e}")
        return None

# --- Additional utility functions for specific file types ---

def get_supported_file_types():
    """Return list of supported file types"""
    return {
        'pdf': ['.pdf'],
        'powerpoint': ['.ppt', '.pptx'],
        'word_document': ['.doc', '.docx'],
        'images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']
    }

def validate_file_type(filepath: str):
    """Validate if file type is supported"""
    supported_types = get_supported_file_types()
    file_extension = os.path.splitext(filepath)[1].lower()
    
    for file_type, extensions in supported_types.items():
        if file_extension in extensions:
            return True, file_type
    
    return False, "unsupported"